# -*- coding: utf-8 -*-
"""生成「研习台」汇报 PPTX（WPS / PowerPoint 可编辑）。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml.etree import SubElement

# ---- 主题色 ----
PAPER  = RGBColor(0xF6,0xF2,0xEA)
INK    = RGBColor(0x2B,0x28,0x23)
INK2   = RGBColor(0x5C,0x55,0x4B)
INK3   = RGBColor(0x8C,0x84,0x75)
CLAY   = RGBColor(0xC2,0x68,0x3B)
CLAY_D = RGBColor(0xA4,0x50,0x2B)
GREEN  = RGBColor(0x3A,0x6B,0x5E)
BLUE   = RGBColor(0x3E,0x6C,0xA8)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
LINE   = RGBColor(0xE2,0xD9,0xC8)
SURF   = RGBColor(0xFF,0xFF,0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0,0,SW,SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s

def bg(s):  # 移除自动加的占位
    return s

def txt(s, x,y,w,h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, italic=False, font="Microsoft YaHei"):
    tb = s.shapes.add_textbox(x,y,w,h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return tb

def bullets(s, x,y,w,h, items, size=16, color=INK, gap=6, marker="•", mcolor=CLAY):
    tb = s.shapes.add_textbox(x,y,w,h); tf = tb.text_frame; tf.word_wrap = True
    for i,it in enumerate(items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.12
        rm = p.add_run(); rm.text = marker+"  "
        rm.font.size = Pt(size); rm.font.color.rgb = mcolor; rm.font.bold = True; rm.font.name="Microsoft YaHei"
        if isinstance(it, tuple):
            head, body = it
            rh = p.add_run(); rh.text = head
            rh.font.size = Pt(size); rh.font.bold = True; rh.font.color.rgb = CLAY_D; rh.font.name="Microsoft YaHei"
            rb = p.add_run(); rb.text = body
            rb.font.size = Pt(size); rb.font.color.rgb = color; rb.font.name="Microsoft YaHei"
        else:
            r = p.add_run(); r.text = it
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name="Microsoft YaHei"
    return tb

def eyebrow(s, x, y, text):
    tb = s.shapes.add_textbox(x,y,Inches(6),Inches(0.35)); tf=tb.text_frame
    p=tf.paragraphs[0]
    r=p.add_run(); r.text=text; r.font.size=Pt(13); r.font.bold=True
    r.font.color.rgb=CLAY_D; r.font.name="Microsoft YaHei"
    # 小横线
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y+Inches(0.34), Inches(0.32), Inches(0.045))
    ln.fill.solid(); ln.fill.fore_color.rgb=CLAY; ln.line.fill.background(); ln.shadow.inherit=False
    return tb

def title(s, x, y, text, size=30, color=INK):
    return txt(s, x, y, Inches(11.5), Inches(0.9), text, size=size, color=color, bold=True)

def card(s, x,y,w,h, tag, head, body, accent=CLAY, head_size=16, body_size=13.5):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x,y,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb = SURF
    sh.line.color.rgb = LINE; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    # 轻微阴影
    el = sh._element.spPr
    import xml.etree.ElementTree as ET
    ef = SubElement(el, qn('a:effectLst'))
    shd = SubElement(ef, qn('a:outerShdw'), {"blur":"90000","dist":"45000","dir":"5400000","rotWithShape":"0"})
    clr = SubElement(shd, qn('a:srgbClr'), {"val":"3C3223"})
    SubElement(clr, qn('a:alpha'), {"val":"22000"})
    # 文本
    tf = sh.text_frame; tf.word_wrap=True
    tf.margin_left=Pt(12); tf.margin_right=Pt(12); tf.margin_top=Pt(10); tf.margin_bottom=Pt(10)
    p=tf.paragraphs[0]
    rt=p.add_run(); rt.text=tag; rt.font.size=Pt(11); rt.font.bold=True; rt.font.color.rgb=accent; rt.font.name="Microsoft YaHei"
    p2=tf.add_paragraph(); p2.space_before=Pt(4)
    rh=p2.add_run(); rh.text=head; rh.font.size=Pt(head_size); rh.font.bold=True; rh.font.color.rgb=INK; rh.font.name="Microsoft YaHei"
    p3=tf.add_paragraph(); p3.space_before=Pt(4)
    rb=p3.add_run(); rb.text=body; rb.font.size=Pt(body_size); rb.font.color.rgb=INK2; rb.font.name="Microsoft YaHei"; rb.font.line_spacing=1.15
    return sh

def box(s, x,y,w,h, text, fill, line=None, tcolor=INK, size=14, bold=True, sub=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x,y,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(1.25)
    sh.shadow.inherit=False
    tf=sh.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf.margin_left=Pt(8); tf.margin_right=Pt(8); tf.margin_top=Pt(4); tf.margin_bottom=Pt(4)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=tcolor; r.font.name="Microsoft YaHei"
    if sub:
        p2=tf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER
        r2=p2.add_run(); r2.text=sub; r2.font.size=Pt(11); r2.font.color.rgb=INK2; r2.font.name="Microsoft YaHei"
    return sh

def arrow(s, x1,y1,x2,y2, color=INK3, w=2.0):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1,y1,x2,y2)
    c.line.color.rgb=color; c.line.width=Pt(w)
    ln=c.line._get_or_add_ln()
    tail=SubElement(ln, qn('a:tailEnd')); tail.set('type','triangle'); tail.set('w','med'); tail.set('len','med')
    head=SubElement(ln, qn('a:headEnd')); head.set('type','none')
    return c

def pagefoot(s, n):
    txt(s, Inches(11.6), Inches(7.05), Inches(1.5), Inches(0.3), f"{n} / 19", size=11, color=INK3, align=PP_ALIGN.RIGHT)

# ============ 1 封面 ============
s = slide()
box(s, Inches(0.7),Inches(0.75),Inches(0.62),Inches(0.62),"研", CLAY, tcolor=WHITE, size=22, bold=True)
txt(s, Inches(1.5),Inches(0.82),Inches(4),Inches(0.5),"研习台  RESEARCH · DESK", size=15, color=INK, bold=True)
txt(s, Inches(1.5),Inches(1.18),Inches(4),Inches(0.3),"科研与生活工作台", size=12, color=INK3)
txt(s, Inches(0.7),Inches(2.3),Inches(12),Inches(1.6),
    "用 AI 重塑开放性知识工作\n——「研习台」科研工作台的设计与实践", size=38, color=INK, bold=True)
txt(s, Inches(0.7),Inches(4.0),Inches(11),Inches(1.0),
    "一个纯前端、跨端同步、由开放 AI 能力驱动的科研与生活工作台：\n把海量顶刊与零散灵感，转化为「可行动的科研线索」。",
    size=17, color=INK2)
# meta
txt(s, Inches(0.7),Inches(5.5),Inches(11.5),Inches(0.5),
    "汇报人：_______      方向：技术经济与管理      日期：2026.09", size=14, color=INK3)
pagefoot(s,1)

# ============ 2 目录 ============
s=slide(); eyebrow(s,Inches(0.7),Inches(0.6),"Agenda"); title(s,Inches(0.7),Inches(0.95),"本次汇报的三条主线")
items=[("01","设计思路与亮点","为什么做、做什么、独特在哪"),
       ("02","技术路线","多端同步实现 + 可复用的工程实践"),
       ("03","其他思考与展望","方法论、元实践、局限与下一步"),
       ("+","演示与 Q&A","结合真实界面走查")]
cw=Inches(5.7); ch=Inches(1.2); gx=Inches(0.5); gy=Inches(0.35)
x0=Inches(0.7); y0=Inches(2.1)
for i,(n,t,d) in enumerate(items):
    r=i//2; c=i%2
    x=x0+c*(cw+gx); y=y0+r*(ch+gy)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,cw,ch)
    sh.fill.solid(); sh.fill.fore_color.rgb=SURF; sh.line.color.rgb=LINE; sh.line.width=Pt(1); sh.shadow.inherit=False
    tf=sh.text_frame; tf.word_wrap=True; tf.margin_left=Pt(14)
    p=tf.paragraphs[0]; rn=p.add_run(); rn.text=n; rn.font.size=Pt(20); rn.font.bold=True; rn.font.color.rgb=CLAY; rn.font.name="Microsoft YaHei"
    rn2=p.add_run(); rn2.text="   "+t; rn2.font.size=Pt(16); rn2.font.bold=True; rn2.font.color.rgb=INK; rn2.font.name="Microsoft YaHei"
    p2=tf.add_paragraph(); p2.space_before=Pt(3)
    rd=p2.add_run(); rd.text="        "+d; rd.font.size=Pt(12.5); rd.font.color.rgb=INK3; rd.font.name="Microsoft YaHei"
pagefoot(s,2)

# ============ 3 背景问题 ============
s=slide(); eyebrow(s,Inches(0.7),Inches(0.55),"Why · 问题"); title(s,Inches(0.7),Inches(0.9),"开放性知识工作的真实困境")
cards=[("信息散落","工具太多，线索太碎","选题、文献、笔记、日程、生活记录分散在 5+ 个 App，缺少统一工作台。",CLAY),
       ("收藏 ≠ 理解","「存了」不等于「读了」","顶刊论文收藏即结束；摘要是英文、方法陌生，阅读门槛高。",GREEN),
       ("跨端割裂","电脑写、手机记，对不上","实验室用电脑、通勤用手机，数据不互通，灵感难落地。",BLUE),
       ("被动接收","信息过载，缺乏过滤","推送 overwhelms；缺少「与我方向相关」的筛选，注意力被稀释。",CLAY_D)]
cw=Inches(5.85); ch=Inches(1.55); gx=Inches(0.4); gy=Inches(0.3); x0=Inches(0.7); y0=Inches(2.0)
for i,(t,h,b,a) in enumerate(cards):
    r=i//2; c=i%2; x=x0+c*(cw+gx); y=y0+r*(ch+gy)
    card(s,x,y,cw,ch,t,h,b,accent=a)
txt(s,Inches(0.7),Inches(6.4),Inches(11.8),Inches(0.6),
    "→ 研习台的目标：把「信息流」变成「可行动的科研线索」，让 AI 替你做过滤、翻译与结构化。",
    size=16, color=CLAY_D, bold=True)
pagefoot(s,3)

# ============ 4 产品定位 ============
s=slide(); eyebrow(s,Inches(0.7),Inches(0.55),"What · 定位"); title(s,Inches(0.7),Inches(0.9),"研习台 = 科研与生活的统一工作台")
txt(s,Inches(0.7),Inches(1.7),Inches(11.8),Inches(0.5),
    "一个打开网页即用的工具：数据自主、AI 辅助、跨端无缝。无后端、无账号、零成本部署。", size=16, color=INK2)
cards=[("🗂 科研主线","顶刊每日推荐、论文追踪（选题→投稿→录用）、项目、方向、研读笔记、成果档案。",CLAY),
       ("🌱 生活侧写","运动、阅读观影、旅行、习惯打卡，以及周/月/年阶段回顾。",GREEN),
       ("⚡ 低摩擦捕获","悬浮按钮 + 底部速记面板，灵感/选题/文献/待办一键入箱。",BLUE)]
cw=Inches(3.85); ch=Inches(1.7); gx=Inches(0.3); x0=Inches(0.7); y0=Inches(2.4)
for i,(h,b,a) in enumerate(cards):
    card(s,x0+i*(cw+gx),y0,cw,ch,"",h,b,accent=a,head_size=16,body_size=13.5)
# flow
box(s,Inches(1.2),Inches(4.7),Inches(3.2),Inches(0.9),"📱 移动端 随手记", SURF, line=LINE, tcolor=INK, size=14)
arrow(s, Inches(4.4),Inches(5.15),Inches(5.4),Inches(5.15), CLAY)
box(s, Inches(5.4),Inches(4.7),Inches(2.6),Inches(0.9),"☁ 云端 单一真相源", GREEN, tcolor=WHITE, size=14)
arrow(s, Inches(8.0),Inches(5.15),Inches(9.0),Inches(5.15), CLAY)
box(s, Inches(9.0),Inches(4.7),Inches(3.2),Inches(0.9),"💻 电脑端 深阅读/管理", SURF, line=LINE, tcolor=INK, size=14)
txt(s, Inches(1.2),Inches(5.8),Inches(11),Inches(0.4),"同一套网页，同一份数据，手机与电脑共享。", size=13, color=INK3, align=PP_ALIGN.CENTER)
pagefoot(s,4)

# ============ 5 设计理念 ============
s=slide(); eyebrow(s,Inches(0.7),Inches(0.55),"Design · 理念"); title(s,Inches(0.7),Inches(0.9),"四条贯穿始终的设计原则")
bullets(s, Inches(0.9),Inches(2.1), Inches(11.6), Inches(4.6), [
  ("本地优先 · 数据自主 —— ","研究者的私密数据默认存本地，可选自建云端；不依赖任何商业平台，随时可导出。"),
  ("AI 做「过滤器 / 翻译器 / 结构器」 —— ","不是堆砌大模型，而是用轻量 API 把信息筛成与你相关的、能读懂的、成体系的科研线索。"),
  ("低摩擦捕获 —— ","灵感稍纵即逝，记录动作要 ≤ 2 次点击；先记下来，再整理。"),
  ("科研与生活一体化 —— ","把生活记录与阶段回顾纳入工作台，让「复盘」成为习惯而非负担。"),
], size=17, gap=14)
pagefoot(s,5)

# ============ 6 亮点一 AI 推荐 ============
s=slide(); eyebrow(s,Inches(0.7),Inches(0.55),"Highlight 01 · AI 核心"); title(s,Inches(0.7),Inches(0.9),"顶刊每日推荐 + 自动相关性打分")
txt(s, Inches(0.7),Inches(1.7),Inches(11.8),Inches(0.5),
    "不是「推论文」，而是按你的研究方向，把海量顶刊变成一份「与你相关的待读清单」。", size=16, color=INK2)
# 流水线图
box(s, Inches(0.7),Inches(2.5),Inches(2.5),Inches(1.1),"OpenAlex\n开放论文库", SURF, line=LINE, tcolor=INK, size=13, sub="免密钥·CORS")
arrow(s, Inches(3.25),Inches(3.05),Inches(3.75),Inches(3.05), INK3)
box(s, Inches(3.75),Inches(2.5),Inches(2.6),Inches(1.1),"UTD24 / FT50\n按 60+ 顶刊过滤", SURF, line=LINE, tcolor=INK, size=13, sub="近 45 天新文")
arrow(s, Inches(6.4),Inches(3.05),Inches(6.9),Inches(3.05), INK3)
box(s, Inches(6.9),Inches(2.5),Inches(2.6),Inches(1.1),"主题/方法词典\n正则匹配打分", SURF, line=LINE, tcolor=INK, size=13, sub="相关度排序")
arrow(s, Inches(9.55),Inches(3.05),Inches(10.1),Inches(3.05), GREEN)
box(s, Inches(10.1),Inches(2.5),Inches(2.5),Inches(1.1),"输出\n待读清单", GREEN, tcolor=WHITE, size=13, sub="自动标注主题/方法")
box(s, Inches(3.75),Inches(4.2),Inches(5.85),Inches(0.9),"「与你方向相关」的待读清单：自动标注主题 / 方法 / 方向关联建议", GREEN, tcolor=WHITE, size=13)
bullets(s, Inches(0.9),Inches(5.4), Inches(11.6), Inches(1.7), [
  ("数据源 OpenAlex： ","免费、免密钥、浏览器可直连，覆盖 UTD24/FT50 等 60+ 顶刊近 45 天新文。"),
  ("方向/方法词典： ","内置数字化转型、AI、创新、专利等主题与 DID/IV/文本分析等方法，自动打分。"),
  ("结果当日缓存： ","每天只请求一次，省流量也更快。"),
], size=13.5, gap=7)
pagefoot(s,6)

# ============ 7 亮点二 翻译 ============
s=slide(); eyebrow(s,Inches(0.7),Inches(0.55),"Highlight 02 · 降门槛"); title(s,Inches(0.7),Inches(0.9),"一键摘要翻译：把顶刊读「懂」")
card(s, Inches(0.7),Inches(2.1),Inches(5.85),Inches(1.7),"主通道","微软 Edge 翻译接口",
     "免密钥、CORS 开放、国内可直连。Token 8 分钟缓存，按论文 id 缓存译文，同一篇只翻一次。", accent=CLAY)
card(s, Inches(6.75),Inches(2.1),Inches(5.85),Inches(1.7),"备用通道","Google gtx 兜底",
     "主通道异常时自动切换，保证「读得了」。带超时控制，单篇不卡死整体。", accent=GREEN)
txt(s, Inches(0.7),Inches(4.2),Inches(11.9),Inches(1.0),
    "价值：把英文摘要、方法说明翻译成中文对照，显著降低顶刊阅读门槛——尤其对快速判断「这篇与我何干」极为高效。",
    size=17, color=INK2)
txt(s, Inches(0.7),Inches(5.5),Inches(11.9),Inches(0.5),
    "设计取舍：用免费公开接口而非自建模型服务，符合「零成本、可直连、易维护」的个人工具定位。",
    size=13.5, color=INK3)
pagefoot(s,7)

# ============ 8 亮点三 全生命周期 ============
s=slide(); eyebrow(s,Inches(0.7),Inches(0.55),"Highlight 03 · 管理"); title(s,Inches(0.7),Inches(0.9),"科研全生命周期，一处管完")
card(s, Inches(0.7),Inches(2.1),Inches(3.85),Inches(2.0),"","📄 论文追踪",
     "选题构思→文献综述→方法设计→初稿→投稿→返修→录用→发表，阶段流一目了然。", accent=CLAY)
card(s, Inches(4.75),Inches(2.1),Inches(3.85),Inches(2.0),"","📁 项目 / 方向",
     "学术研究、学科竞赛、实习实践分类型管理；研究方向独立成册，论文可挂接方向。", accent=GREEN)
card(s, Inches(8.8),Inches(2.1),Inches(3.8),Inches(2.0),"","📝 笔记 / 成果",
     "研读笔记随手写；成果档案（论文/项目/获奖/会议）自动统计，一键导出清单。", accent=BLUE)
txt(s, Inches(0.7),Inches(4.5),Inches(11.9),Inches(1.0),
    "把「想到—读—做—成」的闭环沉淀在同一处，减少在不同工具间搬运的损耗。", size=17, color=INK2)
pagefoot(s,8)

# ============ 9 亮点四 生活回顾 ============
s=slide(); eyebrow(s,Inches(0.7),Inches(0.55),"Highlight 04 · 复盘"); title(s,Inches(0.7),Inches(0.9),"生活侧写 + 阶段回顾")
card(s, Inches(0.7),Inches(2.1),Inches(5.85),Inches(2.2),"","🌿 生活记录",
     "运动（次数/时长/连续天数）、阅读与观影、旅行计划、习惯打卡——科研之外，生活也被温柔记录。", accent=GREEN)
card(s, Inches(6.75),Inches(2.1),Inches(5.85),Inches(2.2),"","📊 周 / 月 / 年回顾",
     "自动生成阶段小结与完成率条形图，把科研进展、生活记录、阅读思考、习惯坚持汇成一份「给自己的复盘」。", accent=CLAY)
txt(s, Inches(0.7),Inches(4.7),Inches(11.9),Inches(1.0),
    "设计主张：研究者首先是完整的人。回顾功能让长期主义可见、可坚持。", size=17, color=GREEN, bold=True)
pagefoot(s,9)

# ============ 10 亮点五 捕获 ============
s=slide(); eyebrow(s,Inches(0.7),Inches(0.55),"Highlight 05 · 捕获"); title(s,Inches(0.7),Inches(0.9),"灵感收集箱 + 快速记录")
bullets(s, Inches(0.9),Inches(2.1), Inches(11.6), Inches(3.5), [
  "右下角悬浮按钮（FAB）随时唤起速记面板，灵感/选题/文献/生活/待办一键分类入箱。",
  "支持「存为今日任务」，捕获即行动，灵感不流失到别处。",
  "首页置顶「随手记一条灵感 / 选题…」，低摩擦、零心智负担，契合知识工作的即兴性。",
], size=17, gap=16)
txt(s, Inches(0.9),Inches(5.4),Inches(11.6),Inches(0.8),
    "这是「低摩擦捕获」原则最直接的产品落点：先收集、后整理，避免好点子沉没。", size=14, color=INK3)
pagefoot(s,10)

# ============ 11 价值小结 ============
s=slide(); eyebrow(s,Inches(0.7),Inches(0.55),"Value · 小结"); title(s,Inches(0.7),Inches(0.9),"研习台带来的改变")
card(s, Inches(0.7),Inches(2.1),Inches(3.85),Inches(2.3),"","从「信息过载」→ 到「相关性筛选」",
     "AI 替你过滤，注意力回到真正相关的文献。", accent=CLAY, head_size=15)
card(s, Inches(4.75),Inches(2.1),Inches(3.85),Inches(2.3),"","从「收藏沉睡」→ 到「读得懂、用得上」",
     "翻译 + 方向关联，让顶刊真正进入思考。", accent=GREEN, head_size=15)
card(s, Inches(8.8),Inches(2.1),Inches(3.8),Inches(2.3),"","从「跨端割裂」→ 到「一处更新、处处同步」",
     "手机记、电脑管，数据自主不锁死。", accent=BLUE, head_size=15)
pagefoot(s,11)

# ============ 12 技术路线总览 ============
s=slide(); eyebrow(s,Inches(0.7),Inches(0.55),"Technical · 总览"); title(s,Inches(0.7),Inches(0.9),"技术路线：一个纯前端的轻量架构")
box(s, Inches(0.7),Inches(2.3),Inches(3.4),Inches(1.4),"浏览器（手机 / 电脑）", SURF, line=LINE, tcolor=INK, size=13, sub="同一套响应式 SPA\nVanilla JS · 零框架")
box(s, Inches(4.5),Inches(2.3),Inches(3.4),Inches(1.4),"静态托管", SURF, line=LINE, tcolor=INK, size=13, sub="GitHub Pages\n一键 deploy.bat 发布")
box(s, Inches(8.3),Inches(1.4),Inches(3.4),Inches(1.2),"开放 API（免密钥）", BLUE, tcolor=WHITE, size=13, sub="OpenAlex · 论文\nEdge/Google · 翻译")
box(s, Inches(8.3),Inches(3.5),Inches(3.4),Inches(1.2),"存储", GREEN, tcolor=WHITE, size=13, sub="localStorage（本地优先）\nJSONBin（可选云端）")
arrow(s, Inches(4.1),Inches(3.0),Inches(4.5),Inches(3.0), INK3)
arrow(s, Inches(7.9),Inches(2.0),Inches(8.3),Inches(1.9), BLUE)
arrow(s, Inches(7.9),Inches(3.4),Inches(8.3),Inches(3.9), GREEN)
txt(s, Inches(0.7),Inches(5.2),Inches(11.9),Inches(1.0),
    "一句话：前端即应用，开放 API 即服务，浏览器即数据库。零服务器成本，个人即可长期维护。", size=17, color=INK2)
pagefoot(s,12)

# ============ 13 核心技术一 多端同步（重点） ============
s=slide(); eyebrow(s,Inches(0.7),Inches(0.55),"Core · 重点"); title(s,Inches(0.7),Inches(0.9),"移动端与电脑端，如何同步？")
txt(s, Inches(0.7),Inches(1.7),Inches(11.9),Inches(0.5),
    "三层方案，层层兜底，保证「一处更新、处处同步」且手机/电脑体验一致。", size=16, color=INK2)
card(s, Inches(0.7),Inches(2.4),Inches(3.85),Inches(2.6),"第一层 · 同源","同一套代码",
     "同一份静态托管（GitHub Pages），手机电脑打开同一 URL 即同一应用。断点 780px 切换导航形态：桌面侧栏 ↔ 移动底部 Tab，数据与逻辑完全共用。", accent=CLAY)
card(s, Inches(4.75),Inches(2.4),Inches(3.85),Inches(2.6),"第二层 · 同步","本地优先 + 云端真相源",
     "整份应用状态 = 一个 JSON。localStorage 本地优先；开启同步后存 JSONBin 云端，启动拉取、本地更旧则采用云端（last-write-wins + 时间戳），保存防抖后推送。", accent=GREEN)
card(s, Inches(8.8),Inches(2.4),Inches(3.8),Inches(2.6),"第三层 · 更新","版本自更新",
     "每次发布更新 version.txt；手机端打开时 checkUpdate 发现新版本自动重载，无需应用商店即可拿到新功能。", accent=BLUE)
txt(s, Inches(0.7),Inches(5.3),Inches(11.9),Inches(0.6),
    "降级策略：未配置云端 / 断网时，自动回退为纯本地（localStorage），功能不中断——「优雅降级」原则的体现。", size=13.5, color=INK3)
pagefoot(s,13)

# ============ 14 同步机制图示 ============
s=slide(); eyebrow(s,Inches(0.7),Inches(0.55),"Core · 机制"); title(s,Inches(0.7),Inches(0.9),"同步时序：手机 ↔ 云端 ↔ 电脑")
box(s, Inches(0.7),Inches(2.2),Inches(3.4),Inches(1.3),"📱 手机端", SURF, line=LINE, tcolor=INK, size=14, sub="localStorage + 推送")
box(s, Inches(4.95),Inches(2.9),Inches(3.4),Inches(1.3),"☁ JSONBin 云端", GREEN, tcolor=WHITE, size=13, sub="单一真相源（一个 Bin）")
box(s, Inches(9.2),Inches(2.2),Inches(3.4),Inches(1.3),"💻 电脑端", SURF, line=LINE, tcolor=INK, size=14, sub="localStorage + 拉取")
arrow(s, Inches(4.95),Inches(3.2),Inches(3.6),Inches(3.0), CLAY)
txt(s, Inches(3.7),Inches(2.6),Inches(1.2),Inches(0.4),"启动 pull", size=11, color=CLAY)
arrow(s, Inches(4.95),Inches(3.5),Inches(9.0),Inches(3.0), CLAY)
txt(s, Inches(8.2),Inches(2.6),Inches(1.2),Inches(0.4),"启动 pull", size=11, color=CLAY)
arrow(s, Inches(3.6),Inches(3.0),Inches(4.95),Inches(3.4), GREEN)
txt(s, Inches(3.9),Inches(3.7),Inches(1.2),Inches(0.4),"保存 push", size=11, color=GREEN)
arrow(s, Inches(9.0),Inches(3.0),Inches(4.95),Inches(3.6), GREEN)
txt(s, Inches(6.6),Inches(3.7),Inches(1.2),Inches(0.4),"保存 push", size=11, color=GREEN)
box(s, Inches(0.7),Inches(5.1),Inches(11.9),Inches(1.4),
    "冲突处理：last-write-wins（以 savedAt 时间戳为准）；推送前 700ms 防抖 + 单飞（in-flight），避免并发覆盖。\n适用场景：单人多端、非并发编辑——个人知识工作的典型形态。",
    SURF, line=LINE, tcolor=INK, size=13)
pagefoot(s,14)

# ============ 15 核心技术二 开放 API ============
s=slide(); eyebrow(s,Inches(0.7),Inches(0.55),"Core · 零成本"); title(s,Inches(0.7),Inches(0.9),"零成本开放 API 架构")
card(s, Inches(0.7),Inches(2.1),Inches(5.85),Inches(1.7),"","📚 OpenAlex（论文）",
     "完全免费、无需密钥、支持浏览器跨域直连。按期刊 ISSN + 时间窗过滤，返回结构化 JSON，前端直接计算相关性。", accent=BLUE)
card(s, Inches(6.75),Inches(2.1),Inches(5.85),Inches(1.7),"","🌐 Edge / Google（翻译）",
     "微软 Edge 翻译接口国内可直连、免密钥；Google gtx 作兜底。两通道皆带超时，保障可用。", accent=CLAY)
title(s, Inches(0.7),Inches(4.1),"可复用的缓存策略", size=20)
bullets(s, Inches(0.9),Inches(4.7), Inches(11.6), Inches(2.0), [
  ("论文推荐： ","按「当天」缓存到 localStorage，每日只请求一次。"),
  ("摘要译文： ","按论文 id 缓存，同一篇只翻译一次；缓存上限 240 条，超量按时间淘汰。"),
  ("统一收益： ","省流量、提速、降低对第三方接口的依赖与受限风险。"),
], size=15, gap=8)
pagefoot(s,15)

# ============ 16 核心技术三 工程实践 ============
s=slide(); eyebrow(s,Inches(0.7),Inches(0.55),"Core · 工程"); title(s,Inches(0.7),Inches(0.9),"值得复用的工程实践")
bullets(s, Inches(0.9),Inches(2.1), Inches(11.6), Inches(4.4), [
  ("单一数据源 + 规范化（normalize）： ","整份应用状态是一个 JSON 对象；读入即归一（补默认字段、修类型），渲染层永远拿到「干净」的结构。"),
  ("版本迁移（v1 → v2）： ","旧数据自动拆分项目/论文、ICS 任务转日历事件，老用户升级无感。"),
  ("渐进增强 / 优雅降级： ","同步失败不阻塞、翻译失败回退、本地存储满则放弃缓存——任何一环出问题，核心功能仍可工作。"),
  ("零框架、零构建： ","纯 HTML/CSS/JS，单文件即可运行；也可打包进 dist/。维护成本低，适合个人长期持有。"),
], size=16, gap=13)
pagefoot(s,16)

# ============ 17 方法论 + 元实践 ============
s=slide(); eyebrow(s,Inches(0.7),Inches(0.55),"Other · 反思"); title(s,Inches(0.7),Inches(0.9),"方法论：AI 在知识工作中的三角色")
card(s, Inches(0.7),Inches(2.1),Inches(3.85),Inches(1.9),"","🔻 过滤器",
     "从信息洪流中筛出「与我相关」，对抗注意力稀释。", accent=CLAY)
card(s, Inches(4.75),Inches(2.1),Inches(3.85),Inches(1.9),"","🔤 翻译器",
     "把语言/方法门槛拉平，让顶刊真正可读可用。", accent=BLUE)
card(s, Inches(8.8),Inches(2.1),Inches(3.8),Inches(1.9),"","🧱 结构器",
     "把零散灵感组织成可追踪的科研进程与复盘。", accent=GREEN)
txt(s, Inches(0.7),Inches(4.4),Inches(11.9),Inches(1.2),
    "元实践视角：这个工具本身的构建，也是「用 AI 辅助完成 AI 工具」的开放知识工作范例——需求、设计、编码、部署一体打通。",
    size=17, color=CLAY_D, bold=True)
pagefoot(s,17)

# ============ 18 局限与下一步 ============
s=slide(); eyebrow(s,Inches(0.7),Inches(0.55),"Other · 展望"); title(s,Inches(0.7),Inches(0.9),"局限与下一步")
card(s, Inches(0.7),Inches(2.1),Inches(5.85),Inches(1.9),"当前局限","同步是「近似一致」",
     "last-write-wins 适合单人多端非并发；若两端同时改，后写覆盖先写，无自动合并。", accent=CLAY)
card(s, Inches(6.75),Inches(2.1),Inches(5.85),Inches(1.9),"下一步","真正的多端实时协同",
     "引入 Yjs / CRDT 做字段级合并，支持离线编辑后无冲突合并；可拓展到多人课题组共用。", accent=GREEN)
card(s, Inches(0.7),Inches(4.2),Inches(5.85),Inches(1.9),"当前局限","翻译依赖第三方免费接口",
     "受配额与可用性影响，需关注稳定性。", accent=CLAY_D)
card(s, Inches(6.75),Inches(4.2),Inches(5.85),Inches(1.9),"下一步","能力延展",
     "RSS 订阅、BibTeX 导出、Notion 同步、AI 摘要总结；并把推荐词典做成可自定义。", accent=BLUE)
pagefoot(s,18)

# ============ 19 总结 ============
s=slide()
box(s, Inches(0.7),Inches(0.8),Inches(0.9),Inches(0.9),"✶", CLAY, tcolor=WHITE, size=34)
eyebrow(s, Inches(1.8),Inches(0.95),"Closing")
title(s, Inches(1.8),Inches(1.3),"让 AI 成为知识工作的「第一道工序」", size=30)
txt(s, Inches(1.8),Inches(2.6),Inches(10.5),Inches(1.8),
    "研习台不是一个更大的收藏夹，而是一套把开放信息转化为个人科研动能的工作方法：同源跨端、数据自主、AI 过滤与结构化。它既是成果，也是方法的载体。",
    size=19, color=INK2)
txt(s, Inches(1.8),Inches(4.8),Inches(10),Inches(0.6),"感谢聆听 · 欢迎讨论与拍砖 🙏", size=16, color=INK3)
pagefoot(s,19)

prs.save(r"E:\工作文件\FHRU_academic_workbuddy\研习台-汇报.pptx")
print("saved", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
